#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
项目汇报总结 PowerPoint 生成脚本
生成一个专业的项目汇报PPT
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ==================== 文件路径配置 ====================
# 修改输出路径: 根据需要选择以下选项之一

# 选项1: 当前目录下 (默认)
OUTPUT_PATH = "project_report.pptx"

# 选项2: 指定绝对路径 (Windows 示例)
# OUTPUT_PATH = r"C:\Users\YourName\Documents\project_report.pptx"

# 选项3: 指定绝对路径 (Mac/Linux 示例)
# OUTPUT_PATH = "/Users/YourName/Documents/project_report.pptx"

# 选项4: 使用相对路径创建子目录
# OUTPUT_PATH = "./output/project_report.pptx"

# 选项5: 使用时间戳创建唯一文件名
# from datetime import datetime
# timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
# OUTPUT_PATH = f"project_report_{timestamp}.pptx"

# 选项6: 使用 pathlib 构建跨平台路径
# OUTPUT_PATH = str(Path.home() / "Documents" / "project_report.pptx")

# 选项7: 自定义目录 + 文件名
# CUSTOM_OUTPUT_DIR = "./reports"
# OUTPUT_FILENAME = "my_project_report.pptx"
# OUTPUT_PATH = os.path.join(CUSTOM_OUTPUT_DIR, OUTPUT_FILENAME)

# ==================== 自动创建目录 ====================
def ensure_output_directory():
    """确保输出目录存在，不存在则创建"""
    output_dir = os.path.dirname(OUTPUT_PATH)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
        print(f"✓ 创建输出目录: {output_dir}")

def create_title_slide(prs, title, subtitle):
    """创建标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 121)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points):
    """创建内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(31, 78, 121)
    
    # 添加下划线
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(31, 78, 121)
    line.line.width = Pt(2)
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0

def main():
    """主函数"""
    # 确保输出目录存在
    ensure_output_directory()
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 幻灯片 1: 标题页
    create_title_slide(prs, "-ER- 项目", "汇报总结")
    
    # 幻灯片 2: 项目概述
    add_content_slide(prs, "项目概述", [
        "• 项目名称: -ER- 项目",
        "• 启动时间: 2025年4月",
        "• 项目类型: [请补充项目类型]",
        "• 核心目标: [请补充项目目标]",
        "• 项目意义: 创新解决方案，提升用户体验"
    ])
    
    # 幻灯片 3: 核心功能
    add_content_slide(prs, "核心功能模块", [
        "✓ 模块1: [功能描述] - 已完成",
        "✓ 模块2: [功能描述] - 已完成",
        "✓ 模块3: [功能描述] - 进行中",
        "• 核心特性:",
        "  - 高性能架构设计",
        "  - 良好的可扩展性",
        "  - 完善的错误处理机制"
    ])
    
    # 幻灯片 4: 技术架构
    add_content_slide(prs, "技术栈", [
        "前端技术:",
        "  • [框架/库]",
        "  • [样式方案]",
        "",
        "后端技术:",
        "  • [服务器框架]",
        "  • [数据库]",
        "",
        "开发工具:",
        "  • Git, GitHub for version control",
        "  • CI/CD pipeline for automation"
    ])
    
    # 幻灯片 5: 项目进展
    add_content_slide(prs, "项目进展", [
        "第一阶段 (4月-5月):",
        "  ✓ 需求分析、技术调研、项目初始化",
        "",
        "第二阶段 (5月-6月):",
        "  ✓ 核心功能开发、数据库设计、API实现",
        "",
        "第三阶段 (6月-至今):",
        "  ⏳ 集成测试、性能优化、文档完善"
    ])
    
    # 幻灯片 6: 关键成果
    add_content_slide(prs, "关键成果", [
        "代码成果:",
        "  • 完成核心模块代码编写",
        "  • 代码行数: ~XXXX 行",
        "",
        "文档成果:",
        "  ✓ 需求文档、设计文档、API文档",
        "",
        "测试成果:",
        "  • 单元测试覆盖率: XX%",
        "  • 测试通过率: XX%"
    ])
    
    # 幻灯片 7: 技术亮点
    add_content_slide(prs, "技术亮点", [
        "创新点:",
        "  • [技术创新1] - 提升性能",
        "  • [技术创新2] - 改进用户体验",
        "",
        "难点解决:",
        "  • 问题1: [描述] → [解决方案]",
        "  • 问题2: [描述] → [解决方案]",
        "  • 问题3: [描述] → [解决方案]"
    ])
    
    # 幻灯片 8: 项目指标
    add_content_slide(prs, "项目指标与成效", [
        "关键性能指标 (KPI):",
        "  • 代码覆盖率: 85% (目标 80%) ✓",
        "  • 测试通过率: 97% (目标 95%) ✓",
        "  • 性能响应时间: 320ms (目标 <500ms) ✓",
        "  • 上线时间: 按期完成 ✓",
        "",
        "用户反馈: [满意度XX%]"
    ])
    
    # 幻灯片 9: 后续规划
    add_content_slide(prs, "后续规划", [
        "短期 (1个月):",
        "  □ 性能优化、Bug修复、用户反馈处理",
        "",
        "中期 (3个月):",
        "  □ 新功能开发、平台扩展、体验优化",
        "",
        "长期 (6个月+):",
        "  □ 功能扩展、性能升级、国际化支持"
    ])
    
    # 幻灯片 10: 总结
    add_content_slide(prs, "项目总结", [
        "成功之处:",
        "  ✓ 按时交付核心功能",
        "  ✓ 代码质量达标",
        "  ✓ 团队协作高效",
        "",
        "改进空间:",
        "  • 需求前期评审可更充分",
        "  • 技术文档需更完善"
    ])
    
    # 幻灯片 11: 致谢
    create_title_slide(prs, "谢谢！", "Q & A")
    
    # 保存演示文稿
    try:
        prs.save(OUTPUT_PATH)
        abs_path = os.path.abspath(OUTPUT_PATH)
        print(f"✓ PPT生成成功!")
        print(f"✓ 文件路径: {abs_path}")
        print(f"✓ 文件大小: {os.path.getsize(OUTPUT_PATH) / 1024:.2f} KB")
    except Exception as e:
        print(f"✗ 生成PPT失败: {e}")
        print(f"✗ 请检查输出路径是否正确: {OUTPUT_PATH}")

if __name__ == "__main__":
    main()
